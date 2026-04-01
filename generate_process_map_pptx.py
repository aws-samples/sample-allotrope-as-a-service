#!/usr/bin/env python3
"""
Generate PowerPoint deck from ASM Service Process Map
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BLUE = RGBColor(0x23, 0x2F, 0x3E)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
LIGHT_BLUE = RGBColor(0x3A, 0x9A, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=WHITE, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
    return shape


# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 1.5, 11, 1.2, "ASM Transformation Service", 44, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11, 0.8, "End-to-End Process Map", 28, LIGHT_BLUE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 11, 0.6, "Converting Laboratory Instrument Data to Allotrope Simple Model (ASM) Format",
             18, DARK_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "Powered by AWS Bedrock Claude  •  Allotropy Library  •  Custom Converters",
             14, ACCENT_BLUE, False, PP_ALIGN.CENTER)


# ── SLIDE 2: Overview ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Service Overview", 32, WHITE, True)
add_text_box(slide, 0.8, 1.3, 11, 0.5,
             "A complete pipeline for converting raw instrument files to validated, integrity-verified ASM output.",
             18, LIGHT_BLUE)

items = [
    "• Hybrid approach: proven allotropy parsers + AI generation for unknown formats",
    "• 31 instruments supported natively via allotropy library",
    "• Custom converter registry for proprietary instrument formats",
    "• AI-powered fallback using AWS Bedrock Claude 3.5 Sonnet",
    "• DVaaS validation with PDF attestation reports",
    "• Data integrity verification: field-by-field source-to-ASM proof",
    "• FDA 21 CFR Part 11 and EMA Annex 11 compliance ready",
    "• No data persistence — files processed in-memory only",
]
add_bullet_slide(slide, 0.8, 2.2, 11.5, 4.5, items, 17)


# ── SLIDE 3: Two Entry Points ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Two Entry Points", 32, WHITE, True)

# Entry Point 1
add_rounded_rect(slide, 0.8, 1.5, 5.5, 1.0, ACCENT_BLUE, "Entry Point 1: CONVERT", 20, WHITE, True)
items1 = [
    '"I have a raw instrument file and need ASM output"',
    "",
    "• Provide: Instrument File + Instrument Config (JSON)",
    "• Full pipeline: Upload → Route → Convert → Validate → Integrity Check",
    "• Returns: ASM file + Validation PDF + Integrity report",
]
add_bullet_slide(slide, 0.8, 2.7, 5.5, 3.5, items1, 15)

# Entry Point 2
add_rounded_rect(slide, 7.0, 1.5, 5.5, 1.0, GREEN, "Entry Point 2: VALIDATE", 20, WHITE, True)
items2 = [
    '"I already have an ASM file and want to check compliance"',
    "",
    "• Provide: ASM JSON file only",
    "• DVaaS validates any ASM file regardless of source",
    "• Returns: Validation status + errors/warnings + PDF attestation",
]
add_bullet_slide(slide, 7.0, 2.7, 5.5, 3.5, items2, 15)


# ── SLIDE 4: Process Flow (5 Steps) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.3, 11, 0.8, "End-to-End Process Flow", 32, WHITE, True)

steps = [
    ("Step 0", "One-Time\nSetup", DARK_GRAY),
    ("Step 1", "Upload", ACCENT_BLUE),
    ("Step 2", "Unified\nConverter", LIGHT_BLUE),
    ("Step 3", "Validation\n(DVaaS)", GREEN),
    ("Step 4", "Data\nIntegrity", ORANGE),
    ("Step 5", "Results", RGBColor(0x8E, 0x44, 0xAD)),
]

x_start = 0.5
for i, (label, name, color) in enumerate(steps):
    x = x_start + i * 2.1
    add_rounded_rect(slide, x, 1.4, 1.8, 0.5, color, label, 13, WHITE, True)
    add_rounded_rect(slide, x, 2.0, 1.8, 0.8, color, name, 14, WHITE, False)

descs = [
    "Create instrument\nconfig JSON\n(once per instrument)",
    "Upload instrument\nfile + config via\nDashboard or API",
    "Intelligent routing:\nCustom → Allotropy\n→ AI fallback",
    "Schema compliance,\nrequired fields,\nPDF attestation",
    "Field-by-field\nsource vs ASM\ncomparison",
    "ASM file, validation\nreport, integrity\nreport, metadata",
]
for i, desc in enumerate(descs):
    x = x_start + i * 2.1
    add_text_box(slide, x, 3.1, 1.8, 1.5, desc, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Arrow indicators
for i in range(5):
    x = x_start + (i + 1) * 2.1 - 0.25
    add_text_box(slide, x, 2.1, 0.3, 0.5, "→", 20, WHITE, True, PP_ALIGN.CENTER)


# ── SLIDE 5: Intelligent Routing (Step 2 Detail) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Step 2: Unified Converter — Intelligent Routing", 28, WHITE, True)
add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "The Unified Converter evaluates routes in priority order and uses the first one that succeeds.",
             16, LIGHT_BLUE)

# Route A
add_rounded_rect(slide, 0.8, 2.0, 3.6, 0.7, ACCENT_BLUE, "Route A: Custom Converter", 16, WHITE, True)
items_a = [
    "• Highest priority",
    "• Exact match: vendor + model",
    "• Registry or embedded converters",
    "• Includes field_mapping ✓",
    "• e.g. Nova FLEX2, EndoScan-V",
]
add_bullet_slide(slide, 0.8, 2.8, 3.6, 3.0, items_a, 13)

# Route B
add_rounded_rect(slide, 4.8, 2.0, 3.6, 0.7, GREEN, "Route B: Allotropy Library", 16, WHITE, True)
items_b = [
    "• Second priority",
    "• 31 supported instruments",
    "• Rule-based parsing",
    "• Field mapping: planned",
    "• e.g. Vi-CELL, SoftMax Pro",
]
add_bullet_slide(slide, 4.8, 2.8, 3.6, 3.0, items_b, 13)

# Route C
add_rounded_rect(slide, 8.8, 2.0, 3.6, 0.7, ORANGE, "Route C: ATaaS (AI)", 16, WHITE, True)
items_c = [
    "• Fallback when A & B fail",
    "• AWS Bedrock Claude 3.5 Sonnet",
    "• Analyzes unknown formats",
    "• Generates Python converter",
    "• Best for simple formats",
]
add_bullet_slide(slide, 8.8, 2.8, 3.6, 3.0, items_c, 13)

# Priority bar
add_text_box(slide, 0.8, 6.2, 11.5, 0.5,
             "Priority:  1. Custom Registry  →  2. Embedded Custom  →  3. Allotropy (31 instruments)  →  4. ATaaS (AI)  →  5. Error",
             14, DARK_GRAY, True, PP_ALIGN.CENTER)


# ── SLIDE 6: Validation (Step 3 Detail) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Step 3: Validation — DVaaS", 32, WHITE, True)
add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "Data Validation as a Service checks ASM structural compliance against Allotrope standards.",
             16, LIGHT_BLUE)

# Checks table
checks = [
    ("Schema compliance", "Does ASM follow Allotrope structure?", "Error"),
    ("Required fields", "manifest, measurement docs, sample docs present?", "Error"),
    ("Measurement IDs", "UUID v4 format?", "Error"),
    ("Timestamps", "ISO 8601 with timezone?", "Error"),
    ("Data source traceability", "Calculated values linked to source?", "Error"),
    ("Unit recognition", "Standard Allotrope units?", "Warning"),
    ("Equipment metadata", "Serial number, software version?", "Warning"),
    ("Device control", "Device type, detection type?", "Warning"),
]

y = 2.0
add_text_box(slide, 0.8, y, 4.0, 0.4, "Check", 14, ACCENT_BLUE, True)
add_text_box(slide, 4.8, y, 5.0, 0.4, "Description", 14, ACCENT_BLUE, True)
add_text_box(slide, 10.0, y, 2.0, 0.4, "Severity", 14, ACCENT_BLUE, True)

for check, desc, severity in checks:
    y += 0.45
    sev_color = RED if severity == "Error" else ORANGE
    add_text_box(slide, 0.8, y, 4.0, 0.4, check, 13, WHITE)
    add_text_box(slide, 4.8, y, 5.0, 0.4, desc, 13, LIGHT_GRAY)
    add_text_box(slide, 10.0, y, 2.0, 0.4, severity, 13, sev_color, True)

# Validation levels
add_text_box(slide, 0.8, 6.0, 11, 0.5, "Validation Levels:", 16, WHITE, True)
add_text_box(slide, 0.8, 6.5, 11, 0.5,
             "Basic (quick structural check)  •  Comprehensive (production use)  •  Certification (regulatory submission)",
             14, LIGHT_GRAY)


# ── SLIDE 7: Data Integrity (Step 4 Detail) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Step 4: Data Integrity Verification", 32, WHITE, True)
add_text_box(slide, 0.8, 1.2, 11, 0.6,
             "Proves every value from the source instrument file was preserved exactly in the ASM output.\n"
             "No rounding, no silent data loss, no transformation.",
             16, LIGHT_BLUE)

items = [
    "• Converter emits field_mapping array during conversion",
    "• Each entry: source field → source value → ASM field → ASM value → unit",
    "• Dashboard compares source_value to asm_value for every entry",
    "• Result: \"100% Data Integrity Confirmed\" or list of mismatches",
]
add_bullet_slide(slide, 0.8, 2.2, 7.0, 2.5, items, 16)

# Implementation status
add_text_box(slide, 0.8, 4.5, 11, 0.5, "Implementation Status:", 18, WHITE, True)

statuses = [
    ("Nova FLEX2 (embedded)", "✓ All 40 CSV columns", GREEN),
    ("Custom Converter Registry", "Spec complete, not yet enforced", ORANGE),
    ("Allotropy Library (31 instruments)", "Planned", DARK_GRAY),
    ("ATaaS (AI)", "Not planned", RED),
]
y = 5.1
for route, status, color in statuses:
    add_text_box(slide, 0.8, y, 5.0, 0.35, route, 14, WHITE)
    add_text_box(slide, 6.0, y, 5.0, 0.35, status, 14, color, True)
    y += 0.4

# Regulatory note
add_text_box(slide, 0.8, 6.8, 11, 0.5,
             "FDA 21 CFR Part 11: proof that electronic records are accurate and unaltered  •  "
             "EMA Annex 11: data integrity throughout the data lifecycle",
             12, DARK_GRAY)


# ── SLIDE 8: Converter Lifecycle ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Converter Lifecycle — Adding New Instruments", 28, WHITE, True)

lifecycle = [
    ("1. IDENTIFY", "Customer has a new instrument\nnot currently supported", DARK_GRAY),
    ("2. CHECK", "Is it in allotropy?\n(31 instruments)", ACCENT_BLUE),
    ("3. BUILD", "Follow Custom Converter\nRequirements doc", LIGHT_BLUE),
    ("4. REGISTER", "POST /register\nStatus: PENDING", ORANGE),
    ("5. APPROVE", "Human reviews code\nPOST /approve", GREEN),
    ("6. AVAILABLE", "Auto-used by Unified\nConverter for matching files", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (label, desc, color) in enumerate(lifecycle):
    x = 0.8 + (i % 3) * 4.1
    y = 1.5 + (i // 3) * 2.8
    add_rounded_rect(slide, x, y, 3.6, 0.6, color, label, 16, WHITE, True)
    add_text_box(slide, x, y + 0.7, 3.6, 1.0, desc, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ── SLIDE 9: Service Architecture ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Service Architecture", 32, WHITE, True)

# Dashboard
add_rounded_rect(slide, 0.8, 1.5, 2.5, 1.0, ACCENT_BLUE, "Dashboard\n(React / CloudFront)", 13, WHITE, True)

# API Gateway
add_rounded_rect(slide, 4.0, 1.5, 3.0, 1.0, LIGHT_BLUE, "API Gateway\n/convert  /validate\n/register  /approve", 12, WHITE, True)

# Lambda Functions
lambdas = [
    ("Unified Converter", 0.5),
    ("Multi-Instrument", 0.5),
    ("ATaaS (AI)", 0.5),
    ("DVaaS", 0.5),
    ("Custom Converter Svc", 0.5),
]
y = 3.0
for name, h in lambdas:
    add_rounded_rect(slide, 4.0, y, 3.0, 0.45, GREEN, name, 12, WHITE, False)
    y += 0.55

# Storage
add_rounded_rect(slide, 8.0, 1.5, 2.5, 1.2, ORANGE, "S3 Buckets\nASM files\nConverters", 13, WHITE, True)
add_rounded_rect(slide, 8.0, 3.0, 2.5, 1.2, ORANGE, "DynamoDB\nConverter Registry\nConversion History", 13, WHITE, True)

# Bedrock
add_rounded_rect(slide, 8.0, 4.8, 2.5, 0.8, RGBColor(0x8E, 0x44, 0xAD),
                 "AWS Bedrock\nClaude 3.5 Sonnet", 13, WHITE, True)

# Labels
add_text_box(slide, 11.0, 1.5, 2.0, 0.4, "AWS Cloud", 14, DARK_GRAY, True)


# ── SLIDE 10: Live Endpoints ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Live Endpoints", 32, WHITE, True)
add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "All services are deployed and operational. Use the Unified Converter for automatic routing.",
             16, LIGHT_BLUE)

endpoints = [
    ("Unified Converter", "tqzatn5bse...amazonaws.com/prod/convert", "Single entry point (recommended)"),
    ("DVaaS", "4ndgjn16zd...amazonaws.com/prod/validate", "ASM validation & certification"),
    ("Custom Converter", "tfv79s08rl...amazonaws.com/prod/", "Register, approve, execute"),
    ("ATaaS", "3dbnsq6w6h...amazonaws.com/prod/convert", "AI-powered conversion (fallback)"),
    ("Multi-Instrument", "6uogqq4zb5...amazonaws.com/prod/convert", "Allotropy library (31 instruments)"),
    ("Dashboard", "d2630v5zyoh8t7.cloudfront.net", "Web UI"),
]

y = 2.2
add_text_box(slide, 0.8, y, 3.0, 0.4, "Service", 14, ACCENT_BLUE, True)
add_text_box(slide, 3.8, y, 5.5, 0.4, "Endpoint", 14, ACCENT_BLUE, True)
add_text_box(slide, 9.5, y, 3.5, 0.4, "Purpose", 14, ACCENT_BLUE, True)

for name, url, purpose in endpoints:
    y += 0.55
    add_text_box(slide, 0.8, y, 3.0, 0.4, name, 14, WHITE, True)
    add_text_box(slide, 3.8, y, 5.5, 0.4, url, 12, LIGHT_GRAY)
    add_text_box(slide, 9.5, y, 3.5, 0.4, purpose, 13, LIGHT_GRAY)


# ── SLIDE 11: Dashboard Tabs ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Dashboard Capabilities", 32, WHITE, True)

tabs = [
    ("Validate ASM File", "Validate an existing ASM file against Allotrope standards", ACCENT_BLUE),
    ("Convert Instrument File", "Convert raw instrument file to validated ASM output", GREEN),
    ("Control Tower", "Monitor conversion jobs, KPIs, and system health", LIGHT_BLUE),
    ("Instrument Config Creator", "Create instrument config JSON files (one-time setup)", ORANGE),
    ("Instrument Registry", "Browse all supported instruments and their status", DARK_GRAY),
    ("Converter Management", "Register, review, and approve custom converters", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (tab, desc, color) in enumerate(tabs):
    x = 0.8 + (i % 2) * 6.2
    y = 1.5 + (i // 2) * 1.8
    add_rounded_rect(slide, x, y, 5.6, 0.6, color, tab, 16, WHITE, True)
    add_text_box(slide, x + 0.2, y + 0.7, 5.2, 0.6, desc, 14, LIGHT_GRAY)


# ── SLIDE 12: FAQ ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Frequently Asked Questions", 32, WHITE, True)

faqs = [
    ("What if my instrument is supported by allotropy but I want field_mapping?",
     "Register a custom converter. Custom converters take priority over allotropy and include field_mapping."),
    ("Do I need a different instrument config for each file?",
     "No. One config per instrument. Reuse it for every file from that instrument."),
    ("Can I validate an ASM file without converting?",
     "Yes. DVaaS validates any ASM file regardless of how it was created."),
    ("What's the difference between DVaaS validation and Data Integrity Verification?",
     "DVaaS checks structural correctness. Integrity verification checks faithfulness to source data."),
    ("How long does conversion take?",
     "Custom converters and allotropy: < 1 second. ATaaS (AI): 5-30 seconds."),
    ("Is my data stored?",
     "Files are processed in-memory only. No data persisted unless explicitly requested."),
]

y = 1.4
for q, a in faqs:
    add_text_box(slide, 0.8, y, 11.5, 0.4, "Q: " + q, 14, ACCENT_BLUE, True)
    add_text_box(slide, 0.8, y + 0.35, 11.5, 0.4, "A: " + a, 13, LIGHT_GRAY)
    y += 0.85


# ── SLIDE 13: Summary ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 1.5, 11, 1.0, "ASM Transformation Service", 40, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11, 0.8, "Production-Ready  •  Regulatory Compliant  •  AI-Powered",
             24, ACCENT_BLUE, False, PP_ALIGN.CENTER)

summary = [
    "✓  31 instruments via allotropy + unlimited via AI",
    "✓  Custom converter registry with human approval workflow",
    "✓  DVaaS validation with PDF attestation reports",
    "✓  Data integrity verification (field-by-field proof)",
    "✓  FDA 21 CFR Part 11 and EMA Annex 11 ready",
    "✓  No data persistence — in-memory processing only",
]
add_bullet_slide(slide, 2.5, 3.8, 8.0, 3.0, summary, 18, WHITE)

add_text_box(slide, 1, 6.5, 11, 0.5,
             "Powered by AWS Lambda  •  AWS Bedrock Claude  •  Allotropy Library  •  AWS CDK",
             14, DARK_GRAY, False, PP_ALIGN.CENTER)


# Save
output_path = r"docs/ASM-Service-Process-Map.pptx"
prs.save(output_path)
print(f"PowerPoint saved to: {output_path}")
